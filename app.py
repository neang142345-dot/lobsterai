import io
import json
import logging
import os
import re
import time
import uuid
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple

import requests
from flask import Flask, jsonify, request, send_file
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas
from docx import Document
from docx.shared import Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# =============================================================================

# App config

# =============================================================================

app = Flask(__name__)
logging.basicConfig(level=os.getenv("LOG_LEVEL", "INFO"))
logger = logging.getLogger("tg_office_assistant")

APP_NAME = os.getenv("APP_NAME", "tg-office-assistant")
APP_ENV = os.getenv("APP_ENV", "production")
HOST = os.getenv("HOST", "0.0.0.0")
PORT = int(os.getenv("PORT", "10000"))
DEBUG = os.getenv("DEBUG", "false").lower() == "true"

BOT_TOKEN = os.getenv("BOT_TOKEN", "")
WEBHOOK_SECRET = os.getenv("WEBHOOK_SECRET", "my-secret")
RENDER_EXTERNAL_URL = os.getenv("RENDER_EXTERNAL_URL", "").rstrip("/")
ADMIN_TOKEN = os.getenv("ADMIN_TOKEN", "")
TELEGRAM_API = f"[https://api.telegram.org/bot{BOT_TOKEN}](https://api.telegram.org/bot{BOT_TOKEN})" if BOT_TOKEN else ""

AI_API_KEY = os.getenv("AI_API_KEY", os.getenv("OPENAI_API_KEY", ""))
AI_API_URL = os.getenv("AI_API_URL", os.getenv("OPENAI_BASE_URL", "[https://api.openai.com/v1/chat/completions](https://api.openai.com/v1/chat/completions)"))
AI_MODEL = os.getenv("AI_MODEL", os.getenv("OPENAI_MODEL", "gpt-4o-mini"))

SEARCH_PROVIDER = os.getenv("SEARCH_PROVIDER", "tavily").lower()
SEARCH_API_KEY = os.getenv("SEARCH_API_KEY", os.getenv("TAVILY_API_KEY", ""))
SEARCH_API_URL = os.getenv("SEARCH_API_URL", "[https://api.tavily.com/search](https://api.tavily.com/search)")
DEFAULT_MAX_RESULTS = int(os.getenv("DEFAULT_MAX_RESULTS", "5"))
DEFAULT_SEARCH_DEPTH = os.getenv("DEFAULT_SEARCH_DEPTH", "advanced")
DEFAULT_TOPIC = os.getenv("DEFAULT_TOPIC", "general")
ENABLE_AUTODETECT_SEARCH = os.getenv("ENABLE_AUTODETECT_SEARCH", "true").lower() == "true"
ENABLE_SAFE_DISCLAIMER = os.getenv("ENABLE_SAFE_DISCLAIMER", "true").lower() == "true"

DEFAULT_LANGUAGE = os.getenv("DEFAULT_LANGUAGE", "zh")
MAX_INPUT_CHARS = int(os.getenv("MAX_INPUT_CHARS", "8000"))
HTTP_TIMEOUT = int(os.getenv("HTTP_TIMEOUT", "45"))
TEMP_DIR = os.getenv("TEMP_DIR", "/tmp")
os.makedirs(TEMP_DIR, exist_ok=True)

TASKS: Dict[str, Dict[str, Any]] = {}

# =============================================================================

# Helpers

# =============================================================================

def now_iso() -> str:
return datetime.utcnow().isoformat() + "Z"

def update_task(task_id: str, **kwargs: Any) -> None:
task = TASKS.setdefault(task_id, {"task_id": task_id, "created_at": now_iso()})
task.update(kwargs)
task["updated_at"] = now_iso()

def normalize_text(text: str) -> str:
if not text:
return ""
return text.replace("\r\n", "\n").strip()

def clean_text(value: str) -> str:
value = value or ""
value = re.sub(r"\s+", " ", value).strip()
return value

def sanitize_filename(name: str) -> str:
if not name:
return "文件"
name = name.strip()
name = re.sub(r'[\/:*?"<>|]+', "-", name)
name = re.sub(r"\s+", " ", name).strip()
return (name or "文件")[:80]

def detect_language(text: str) -> str:
text = text or ""
zh_chars = len(re.findall(r"[\u4e00-\u9fff]", text))
latin_chars = len(re.findall(r"[A-Za-z]", text))
if zh_chars >= max(5, latin_chars):
return "zh"
if latin_chars > 0:
return "en"
return DEFAULT_LANGUAGE

def make_error(message: str, code: int = 400, extra: Optional[Dict[str, Any]] = None):
payload = {"ok": False, "error": message}
if extra:
payload.update(extra)
return jsonify(payload), code

def build_title(text: str) -> str:
base = normalize_text(text)
if not base:
return "自动生成文档"
first_line = base.splitlines()[0].strip()
return first_line if first_line else "自动生成文档"

def build_output_filename(content: str, ext: str, custom_name: str = None, ai_title: str = None) -> str:
if custom_name:
base_name = sanitize_filename(custom_name)
elif ai_title:
base_name = sanitize_filename(ai_title)
else:
base_name = sanitize_filename(build_title(content))
timestamp = int(time.time())
return os.path.join(TEMP_DIR, f"{base_name}_{timestamp}.{ext}")

# =============================================================================

# Telegram

# =============================================================================

def send_message(chat_id: int, text: str):
if not TELEGRAM_API:
raise RuntimeError("BOT_TOKEN not configured")
resp = requests.post(
f"{TELEGRAM_API}/sendMessage",
json={"chat_id": chat_id, "text": text},
timeout=30,
)
resp.raise_for_status()
data = resp.json()
if not data.get("ok"):
raise RuntimeError(f"Telegram sendMessage failed: {data}")
return data

def send_document(chat_id: int, file_path: str, caption: str = "文件已生成"):
if not TELEGRAM_API:
raise RuntimeError("BOT_TOKEN not configured")
with open(file_path, "rb") as f:
resp = requests.post(
f"{TELEGRAM_API}/sendDocument",
data={"chat_id": chat_id, "caption": caption},
files={"document": f},
timeout=120,
)
resp.raise_for_status()
data = resp.json()
if not data.get("ok"):
raise RuntimeError(f"Telegram sendDocument failed: {data}")
return data

# =============================================================================

# Input extraction

# =============================================================================

def extract_custom_filename(text: str) -> Tuple[Optional[str], str]:
lines = normalize_text(text).splitlines()
custom_name = None
remaining_lines = []

```
for line in lines:
    stripped = line.strip()
    if custom_name is None and (
        stripped.startswith("文件名:") or stripped.startswith("文件名：")
    ):
        custom_name = stripped.split(":", 1)[1].strip() if ":" in stripped else stripped.split("：", 1)[1].strip()
    else:
        remaining_lines.append(line)

return custom_name, "\n".join(remaining_lines).strip()
```

def extract_style(text: str) -> Tuple[Optional[str], str]:
lines = normalize_text(text).splitlines()
style = None
remaining_lines = []

```
for line in lines:
    stripped = line.strip()
    if style is None and (stripped.startswith("风格:") or stripped.startswith("风格：")):
        style = stripped.split(":", 1)[1].strip() if ":" in stripped else stripped.split("：", 1)[1].strip()
    else:
        remaining_lines.append(line)

return style, "\n".join(remaining_lines).strip()
```

def extract_theme(text: str) -> Tuple[Optional[str], str]:
lines = normalize_text(text).splitlines()
theme = None
remaining_lines = []

```
for line in lines:
    stripped = line.strip()
    if theme is None and (stripped.startswith("主题:") or stripped.startswith("主题：")):
        theme = stripped.split(":", 1)[1].strip() if ":" in stripped else stripped.split("：", 1)[1].strip()
    else:
        remaining_lines.append(line)

return theme, "\n".join(remaining_lines).strip()
```

def strip_control_phrases(text: str) -> str:
cleaned = normalize_text(text)
prefixes = [
"生成 pdf:", "生成pdf:", "生成 pdf：", "生成pdf：",
"生成 word:", "生成word:", "生成 word：", "生成word：",
"生成 excel:", "生成excel:", "生成 excel：", "生成excel：",
"生成 ppt:", "生成ppt:", "生成 ppt：", "生成ppt：",
"生成 docx:", "生成docx:", "生成 docx：", "生成docx：",
"生成 xlsx:", "生成xlsx:", "生成 xlsx：", "生成xlsx：",
"生成 pptx:", "生成pptx:", "生成 pptx：", "生成pptx：",
"制作 pdf:", "制作pdf:", "制作 pdf：", "制作pdf：",
"制作 word:", "制作word:", "制作 word：", "制作word：",
"制作 excel:", "制作excel:", "制作 excel：", "制作excel：",
"制作 ppt:", "制作ppt:", "制作 ppt：", "制作ppt：",
"帮我生成", "帮我制作", "生成文件", "生成文档", "制作文件", "制作文档", "制作",
"发我pdf", "回传pdf", "发我word", "发我excel", "发我ppt", "整理成pdf", "整理成word", "整理成ppt",
]

```
lowered = cleaned.lower()
for prefix in prefixes:
    if lowered.startswith(prefix.lower()):
        cleaned = cleaned[len(prefix):].strip()
        break
return cleaned
```

def extract_effective_request(text: str, reply_to_message: dict) -> str:
current_text = normalize_text(text)
reply_text = ""

```
if reply_to_message:
    reply_text = normalize_text(reply_to_message.get("text", ""))

command_like = any(kw in current_text.lower() for kw in [
    "pdf", "docx", "word", "xlsx", "excel", "ppt", "pptx", "生成", "做好", "发我", "整理"
])

if reply_text and command_like and len(current_text) <= 30:
    return reply_text + "\n" + current_text

return current_text
```

# =============================================================================

# Theme system

# =============================================================================

def get_theme_palette(theme: str):
theme = (theme or "商务蓝").strip()

```
if theme == "深色商务":
    return {
        "name": "深色商务",
        "primary": RGBColor(25, 42, 86),
        "secondary": RGBColor(46, 64, 113),
        "accent": RGBColor(93, 173, 226),
        "text": RGBColor(255, 255, 255),
        "subtext": RGBColor(220, 220, 220),
        "light_fill": "D6EAF8",
        "excel_header": "2E4053",
    }

if theme == "极简白":
    return {
        "name": "极简白",
        "primary": RGBColor(40, 40, 40),
        "secondary": RGBColor(150, 150, 150),
        "accent": RGBColor(91, 155, 213),
        "text": RGBColor(40, 40, 40),
        "subtext": RGBColor(120, 120, 120),
        "light_fill": "F2F4F4",
        "excel_header": "DDEBF7",
    }

return {
    "name": "商务蓝",
    "primary": RGBColor(44, 95, 153),
    "secondary": RGBColor(91, 155, 213),
    "accent": RGBColor(23, 162, 184),
    "text": RGBColor(30, 30, 30),
    "subtext": RGBColor(100, 100, 100),
    "light_fill": "D9EAF7",
    "excel_header": "D9EAF7",
}
```

def choose_subtitle_by_style(style: str):
if style == "汇报":
return "管理层汇报材料"
if style == "商务":
return "商务演示文稿"
if style == "正式":
return "正式演示文稿"
return "业务分析报告"

# =============================================================================

# Search abstraction

# =============================================================================

class SearchProvider:
def search(self, query: str, language: str, max_results: int, topic: str, search_depth: str) -> Dict[str, Any]:
raise NotImplementedError

class TavilySearchProvider(SearchProvider):
def **init**(self, api_key: str, api_url: str):
self.api_key = api_key
self.api_url = api_url

```
def search(self, query: str, language: str, max_results: int, topic: str, search_depth: str) -> Dict[str, Any]:
    if not self.api_key:
        raise ValueError("Missing SEARCH_API_KEY or TAVILY_API_KEY")

    payload = {
        "query": query,
        "search_depth": search_depth,
        "max_results": max_results,
        "topic": topic,
        "include_answer": False,
        "include_raw_content": False,
        "include_images": False,
    }
    headers = {
        "Authorization": f"Bearer {self.api_key}",
        "Content-Type": "application/json",
    }
    resp = requests.post(self.api_url, headers=headers, json=payload, timeout=HTTP_TIMEOUT)
    resp.raise_for_status()
    data = resp.json()
    normalized = []
    for item in data.get("results", []) if isinstance(data, dict) else []:
        normalized.append(
            {
                "title": item.get("title", ""),
                "url": item.get("url", ""),
                "content": clean_text(item.get("content", "")),
                "score": item.get("score"),
                "published_date": item.get("published_date"),
            }
        )
    return {
        "provider": "tavily",
        "query": query,
        "language": language,
        "results": normalized,
        "raw": data,
    }
```

class MockSearchProvider(SearchProvider):
def search(self, query: str, language: str, max_results: int, topic: str, search_depth: str) -> Dict[str, Any]:
return {
"provider": "mock",
"query": query,
"language": language,
"results": [
{
"title": f"Mock result for: {query}",
"url": "[https://example.com/mock](https://example.com/mock)",
"content": f"This is a mock search result for query: {query}",
"score": 1.0,
"published_date": None,
}
],
"raw": {},
}

def get_search_provider() -> SearchProvider:
if SEARCH_PROVIDER == "tavily":
return TavilySearchProvider(SEARCH_API_KEY, SEARCH_API_URL)
return MockSearchProvider()

def fuse_search_results(search_data: Dict[str, Any], max_items: int = 5) -> Dict[str, Any]:
results = search_data.get("results", [])[:max_items]
summaries = []
sources = []
for idx, item in enumerate(results, start=1):
title = clean_text(item.get("title", ""))
content = clean_text(item.get("content", ""))
url = item.get("url", "")
summaries.append(f"[{idx}] {title}: {content[:500]}")
sources.append({"index": idx, "title": title, "url": url})
return {
"summary_text": "\n".join(summaries),
"sources": sources,
}

# =============================================================================

# Schema system

# =============================================================================

SCHEMAS: Dict[str, Dict[str, Any]] = {
"work_certificate": {
"display_name": {"zh": "工作证明", "en": "Employment Certificate"},
"fields": [
"title", "employee_name", "employee_id", "company_name", "department", "position",
"employment_start_date", "employment_end_date", "purpose", "issuer", "issue_date",
"contact_info", "notes"
],
},
"contract": {
"display_name": {"zh": "合同/协议", "en": "Contract / Agreement"},
"fields": [
"title", "party_a", "party_b", "background", "scope", "deliverables", "payment_terms",
"timeline", "acceptance", "confidentiality", "liability", "termination", "governing_law",
"signing_date", "notes"
],
},
"company_policy": {
"display_name": {"zh": "公司制度", "en": "Company Policy"},
"fields": [
"title", "purpose", "scope", "roles_and_responsibilities", "policy_rules", "process",
"exceptions", "compliance", "effective_date", "revision_history"
],
},
"invitation_letter": {
"display_name": {"zh": "邀请函", "en": "Invitation Letter"},
"fields": [
"title", "inviter", "invitee", "event_name", "event_purpose", "event_time",
"event_location", "agenda", "contact_info", "closing", "issue_date"
],
},
"notice": {
"display_name": {"zh": "通知", "en": "Notice"},
"fields": [
"title", "audience", "background", "key_message", "effective_time", "actions_required",
"deadline", "contact_info", "issuer", "issue_date"
],
},
"meeting_minutes": {
"display_name": {"zh": "会议纪要", "en": "Meeting Minutes"},
"fields": [
"title", "meeting_time", "meeting_location", "participants", "agenda", "discussion_summary",
"decisions", "action_items", "next_meeting", "recorder"
],
},
"application": {
"display_name": {"zh": "申请书", "en": "Application Letter"},
"fields": [
"title", "applicant", "recipient", "purpose", "background", "requested_action",
"supporting_points", "attachments", "closing", "date"
],
},
"industry_analysis": {
"display_name": {"zh": "行业分析", "en": "Industry Analysis"},
"fields": [
"title", "industry_scope", "executive_summary", "market_size", "growth_drivers",
"competitive_landscape", "risks", "policy_factors", "outlook", "sources_summary"
],
},
"report": {
"display_name": {"zh": "汇报材料", "en": "Report"},
"fields": [
"title", "summary", "background", "current_status", "key_issues", "progress",
"data_points", "next_steps", "support_needed", "appendix"
],
},
"table_spec": {
"display_name": {"zh": "表格字段规范", "en": "Table Field Specification"},
"fields": [
"title", "table_name", "business_purpose", "fields", "validation_rules", "indexing",
"permissions", "examples", "change_log"
],
},
"salary_table": {
"display_name": {"zh": "工资表", "en": "Salary Table"},
"fields": ["title", "columns", "rows", "notes"],
},
"sales_table": {
"display_name": {"zh": "销售表", "en": "Sales Table"},
"fields": ["title", "columns", "rows", "notes"],
},
"presentation": {
"display_name": {"zh": "演示文稿", "en": "Presentation"},
"fields": ["title", "summary", "slides", "notes"],
},
"generic_document": {
"display_name": {"zh": "通用文档", "en": "Generic Document"},
"fields": ["title", "summary", "sections", "notes"],
},
}

KEYWORD_TO_SCHEMA = {
"工作证明": "work_certificate",
"证明": "work_certificate",
"合同": "contract",
"协议": "contract",
"制度": "company_policy",
"政策": "company_policy",
"邀请函": "invitation_letter",
"通知": "notice",
"会议纪要": "meeting_minutes",
"纪要": "meeting_minutes",
"申请书": "application",
"请假申请": "application",
"行业分析": "industry_analysis",
"市场分析": "industry_analysis",
"公司经营分析": "industry_analysis",
"汇报": "report",
"报告": "report",
"表格字段规范": "table_spec",
"字段规范": "table_spec",
"工资表": "salary_table",
"销售表": "sales_table",
"统计表": "sales_table",
"PPT": "presentation",
"ppt": "presentation",
"老板汇报": "presentation",
}

SENSITIVE_SCHEMAS = {"contract", "company_policy", "work_certificate"}

# =============================================================================

# Local fallback intent route

# =============================================================================

def detect_intent_and_file_type_fallback(text: str):
raw = text or ""
t = raw.lower()

```
if "pdf" in t:
    return {"file_type": "pdf", "intent": "generic", "schema_name": "generic_document"}
if "docx" in t or "word" in t:
    return {"file_type": "docx", "intent": "document", "schema_name": "generic_document"}
if "xlsx" in t or "excel" in t:
    return {"file_type": "xlsx", "intent": "table", "schema_name": "sales_table"}
if "pptx" in t or re.search(r"\bppt\b", t):
    return {"file_type": "pptx", "intent": "presentation", "schema_name": "presentation"}

if "表格" in raw or "工资表" in raw or "销售表" in raw or "统计表" in raw or "清单" in raw:
    if "工资" in raw:
        return {"file_type": "xlsx", "intent": "salary_table", "schema_name": "salary_table"}
    if "销售" in raw:
        return {"file_type": "xlsx", "intent": "sales_table", "schema_name": "sales_table"}
    return {"file_type": "xlsx", "intent": "table", "schema_name": "sales_table"}

if "汇报" in raw or "演示" in raw or "路演" in raw or "老板汇报" in raw or "公司经营分析" in raw:
    return {"file_type": "pptx", "intent": "presentation", "schema_name": "presentation"}

if "请假申请" in raw or "申请书" in raw:
    return {"file_type": "docx", "intent": "application", "schema_name": "application"}

if "合同" in raw or "协议" in raw:
    return {"file_type": "docx", "intent": "contract", "schema_name": "contract"}

if "会议纪要" in raw:
    return {"file_type": "docx", "intent": "meeting_minutes", "schema_name": "meeting_minutes"}

if "周报" in raw or "月报" in raw or "日报" in raw or "总结" in raw:
    return {"file_type": "docx", "intent": "weekly_report", "schema_name": "report"}

return {"file_type": "pdf", "intent": "generic", "schema_name": "generic_document"}
```

def enhance_content(intent: str, content: str, style: str = None) -> str:
clean = content.strip() if content and content.strip() else "（空内容）"
lines = [line.strip() for line in clean.splitlines() if line.strip()]

```
if intent == "application":
    if len(lines) <= 3:
        clean = (
            "申请书\n\n"
            "尊敬的领导：\n"
            f"本人因{clean}，特申请相关安排，请予批准。\n\n"
            "此致\n"
            "敬礼"
        )

elif intent == "meeting_minutes":
    clean = (
        "会议纪要\n\n"
        f"{clean}\n\n"
        "一、会议结论\n"
        "待补充\n\n"
        "二、行动项\n"
        "1. 待补充\n"
        "2. 待补充"
    )

elif intent == "weekly_report":
    clean = (
        "工作汇报\n\n"
        "一、本期完成\n"
        f"{clean}\n\n"
        "二、存在问题\n"
        "待补充\n\n"
        "三、下一步计划\n"
        "待补充"
    )

elif intent == "contract":
    if len(lines) <= 5:
        clean = (
            "合作协议\n\n"
            "甲方：\n"
            "乙方：\n"
            "合作背景：\n"
            f"{clean}\n\n"
            "合作内容：\n"
            "付款方式：\n"
            "双方权责：\n"
            "违约责任：\n"
        )

return clean
```

# =============================================================================

# AI / pipeline

# =============================================================================

def infer_schema_type(task: str, details: str) -> str:
haystack = f"{task or ''} {details or ''}".lower()
for keyword, schema_name in KEYWORD_TO_SCHEMA.items():
if keyword.lower() in haystack:
return schema_name
return "generic_document"

def infer_need_search(task: str, details: str, explicit_need_search: Optional[bool]) -> bool:
if explicit_need_search is not None:
return explicit_need_search
if not ENABLE_AUTODETECT_SEARCH:
return False

```
text = f"{task or ''} {details or ''}"
triggers = [
    "完整版", "模板", "范本", "外部资料", "公开资料", "网页", "网络", "政策",
    "行业", "市场", "公司制度", "合同", "协议", "邀请函", "分析", "法规",
    "public", "web", "search", "template", "sample", "policy", "market", "industry",
]
return any(t.lower() in text.lower() for t in triggers)
```

def rewrite_search_query(task: str, details: str, schema_name: str, language: str) -> str:
schema_label = SCHEMAS.get(schema_name, {}).get("display_name", {}).get(language, schema_name)
base = clean_text(f"{task} {details}")
if language == "zh":
return clean_text(f"{schema_label} 完整范本 公开资料 结构要求 示例 {base}")
return clean_text(f"{schema_label} complete sample public references structure example {base}")

def build_schema_prompt(schema_name: str, language: str) -> str:
schema = SCHEMAS.get(schema_name, SCHEMAS["generic_document"])
fields = schema["fields"]
label = schema["display_name"].get(language, schema_name)
return (
f"Document type: {label}\n"
f"Required fields: {', '.join(fields)}\n"
"Return strict JSON only.\n"
"Keys: file_type, intent, schema_name, title, style, summary, content_text, table_headers, table_rows, slides, warnings.\n"
"slides must be an array of {title, bullets}.\n"
"table_headers and table_rows are used for xlsx generation.\n"
)

def call_ai(prompt: str) -> Optional[str]:
if not AI_API_KEY or not AI_API_URL:
return None

```
try:
    resp = requests.post(
        AI_API_URL,
        headers={
            "Authorization": f"Bearer {AI_API_KEY}",
            "Content-Type": "application/json",
        },
        json={
            "model": AI_MODEL,
            "messages": [
                {
                    "role": "system",
                    "content": (
                        "你是一个专业办公文件助手。"
                        "你必须把用户需求转换成最终可生成文件的结构化结果。"
                        "不要输出额外解释，不要使用 markdown 代码块。"
                        "你的目标是帮助系统生成真实文件：pdf、docx、xlsx、pptx。"
                    ),
                },
                {"role": "user", "content": prompt},
            ],
            "temperature": 0.3,
            "response_format": {"type": "json_object"},
        },
        timeout=HTTP_TIMEOUT,
    )
    if resp.status_code != 200:
        logger.warning("AI call failed status=%s body=%s", resp.status_code, resp.text[:500])
        return None
    data = resp.json()
    return data["choices"][0]["message"]["content"]
except Exception:
    logger.exception("call_ai failed")
    return None
```

def normalize_ai_result(data: dict, fallback: dict = None) -> dict:
fallback = fallback or {}
result = {
"file_type": data.get("file_type") or fallback.get("file_type", "pdf"),
"intent": data.get("intent") or fallback.get("intent", "generic"),
"schema_name": data.get("schema_name") or fallback.get("schema_name", "generic_document"),
"title": data.get("title") or fallback.get("title", "文件"),
"style": data.get("style") or fallback.get("style", "默认"),
"summary": data.get("summary") or fallback.get("summary", ""),
"content_text": data.get("content_text") or fallback.get("content_text", ""),
"table_headers": data.get("table_headers") or fallback.get("table_headers", []) or [],
"table_rows": data.get("table_rows") or fallback.get("table_rows", []) or [],
"slides": data.get("slides") or fallback.get("slides", []) or [],
"warnings": data.get("warnings") or fallback.get("warnings", []) or [],
}

```
if result["file_type"] not in {"pdf", "docx", "xlsx", "pptx"}:
    result["file_type"] = fallback.get("file_type", "pdf")

if not isinstance(result["table_headers"], list):
    result["table_headers"] = []
if not isinstance(result["table_rows"], list):
    result["table_rows"] = []
if not isinstance(result["slides"], list):
    result["slides"] = []
if not isinstance(result["warnings"], list):
    result["warnings"] = [str(result["warnings"])]

return result
```

def ai_analyze(
text: str,
style: str = None,
theme: str = None,
task: str = None,
language: str = "zh",
schema_name: str = "generic_document",
need_search: bool = False,
fused_search: Optional[Dict[str, Any]] = None,
forced_file_type: Optional[str] = None,
) -> Optional[dict]:
style_text = style if style else "默认"
search_summary = fused_search.get("summary_text", "") if fused_search else ""

```
prompt = f"""
```

请分析用户需求，并返回严格 JSON，不要输出任何额外解释，不要使用 markdown 代码块。

用户输入:
{text}

任务名:
{task or ''}

文档语言:
{language}

文档风格:
{style_text}

视觉主题:
{theme or '商务蓝'}

schema:
{schema_name}

是否启用网页搜索:
{need_search}

搜索资料摘要:
{search_summary or 'N/A'}

{build_schema_prompt(schema_name, language)}

规则:

1. 目标是生成真实文件，而不是解释说明。
2. 表格、清单、统计、工资表、销售表 -> xlsx。
3. 汇报、演示、路演、PPT、老板汇报、公司经营分析 -> pptx。
4. 请假申请、合同、会议纪要、周报、月报、日报、正式文书 -> docx。
5. 其他一般内容 -> pdf 或 docx。
6. title 必须简洁明确。
7. style 使用: 正式 / 商务 / 汇报 / 默认。
8. content_text 必须是完整可直接写入文件的内容。
9. 如果 file_type 不是 xlsx，table_headers 和 table_rows 返回空数组。
10. 如果 file_type 不是 pptx，slides 返回空数组。
11. 如果 file_type 是 xlsx，尽量生成完整表头，并至少生成 3 行合理示例数据。
12. 如果 file_type 是 pptx:

* 尽量拆成 4 到 6 页
* 每页只保留核心信息
* 每页最多 4 个 bullet
* bullet 必须是短句，不要写成长段

13. 如果是合同、制度、证明类，warnings 中加入“仅供模板/草稿参考”。
14. 如果有搜索摘要，优先吸收其结构与常见字段。
15. 如果明确指定输出格式，则尊重该格式：{forced_file_type or '未指定'}。
    """
    result = call_ai(prompt)
    if not result:
    return None

    try:
    return json.loads(result)
    except Exception:
    logger.exception("AI result json parse failed")
    return None

def build_fallback_result(
cleaned: str,
style: Optional[str],
schema_name: str,
fallback_route: Dict[str, Any],
forced_type: Optional[str],
fused_search: Optional[Dict[str, Any]],
) -> dict:
intent = fallback_route["intent"]
file_type = forced_type or fallback_route["file_type"]
content_text = enhance_content(intent, cleaned, style)

```
if fused_search and fused_search.get("summary_text") and file_type in {"docx", "pdf"}:
    content_text = content_text + "\n\n参考资料摘要：\n" + fused_search["summary_text"]

result = {
    "file_type": file_type,
    "intent": intent,
    "schema_name": schema_name,
    "title": build_title(cleaned),
    "style": style or "默认",
    "summary": "",
    "content_text": content_text,
    "table_headers": [],
    "table_rows": [],
    "slides": [],
    "warnings": [],
}

if file_type == "xlsx":
    if "工资" in cleaned:
        result["table_headers"] = ["姓名", "部门", "基本工资", "绩效", "应发工资"]
        result["table_rows"] = [
            ["张三", "销售部", "8000", "1500", "9500"],
            ["李四", "市场部", "9000", "1200", "10200"],
            ["王五", "运营部", "8500", "1000", "9500"],
        ]
    else:
        result["table_headers"] = ["列1", "列2", "列3", "列4"]
        result["table_rows"] = [
            ["值1", "值2", "值3", "值4"],
            ["值1", "值2", "值3", "值4"],
            ["值1", "值2", "值3", "值4"],
        ]

if file_type == "pptx":
    bullets = [line.strip() for line in cleaned.splitlines() if line.strip()][:12]
    chunks = [bullets[i:i + 4] for i in range(0, len(bullets), 4)] or [["（空内容）"]]
    result["slides"] = []
    for idx, chunk in enumerate(chunks[:6], start=1):
        result["slides"].append({"title": f"第{idx}部分", "bullets": chunk})

if ENABLE_SAFE_DISCLAIMER and schema_name in SENSITIVE_SCHEMAS:
    result["warnings"].append("本内容仅供模板/草稿参考，请结合实际业务、地区法规或专业意见后正式使用。")
    if file_type in {"docx", "pdf"}:
        result["content_text"] += "\n\n风险提示：本内容仅供模板/草稿参考，请结合实际业务、地区法规或专业意见后正式使用。"

return result
```

def create_generation_result(payload: Dict[str, Any], task_id: str) -> Dict[str, Any]:
start_ts = time.time()
update_task(task_id, status="running", step="received", request=payload)

```
task = clean_text(payload.get("task", ""))
details = normalize_text(payload.get("details", payload.get("content", "")))
source_text = normalize_text(payload.get("source_text", ""))
merged_text = normalize_text((task + "\n" + details + "\n" + source_text).strip())
language = payload.get("language") or detect_language(merged_text)
style = payload.get("style") or "默认"
theme = payload.get("theme") or "商务蓝"
custom_name = payload.get("filename")

explicit_need_search = payload.get("need_search") if "need_search" in payload else None
if isinstance(explicit_need_search, str):
    explicit_need_search = explicit_need_search.lower() in {"true", "1", "yes", "y"}

forced_file_type = payload.get("format")
if forced_file_type in {"word", "doc"}:
    forced_file_type = "docx"
if forced_file_type in {"excel", "xls"}:
    forced_file_type = "xlsx"
if forced_file_type in {"ppt"}:
    forced_file_type = "pptx"
if forced_file_type and forced_file_type not in {"pdf", "docx", "xlsx", "pptx"}:
    forced_file_type = None

schema_name = payload.get("schema_name") or infer_schema_type(task, merged_text)
if schema_name not in SCHEMAS:
    schema_name = "generic_document"
update_task(task_id, step="schema_inferred", schema_name=schema_name, language=language)

need_search = infer_need_search(task, merged_text, explicit_need_search)
search_data = None
fused_search = None
if need_search:
    try:
        search_query = payload.get("search_query") or rewrite_search_query(task, merged_text, schema_name, language)
        update_task(task_id, step="searching", search_query=search_query)
        provider = get_search_provider()
        search_data = provider.search(
            query=search_query,
            language=language,
            max_results=int(payload.get("max_results", DEFAULT_MAX_RESULTS)),
            topic=payload.get("topic", DEFAULT_TOPIC),
            search_depth=payload.get("search_depth", DEFAULT_SEARCH_DEPTH),
        )
        fused_search = fuse_search_results(search_data, max_items=int(payload.get("max_results", DEFAULT_MAX_RESULTS)))
        update_task(task_id, step="search_completed", search_sources=fused_search.get("sources", []))
    except Exception as exc:
        logger.exception("search failed, fallback to no-search pipeline")
        update_task(task_id, step="search_failed", search_error=str(exc))
else:
    update_task(task_id, step="search_skipped")

fallback_route = detect_intent_and_file_type_fallback(merged_text)
fallback_route["schema_name"] = schema_name

update_task(task_id, step="ai_generating")
ai_raw = ai_analyze(
    text=merged_text,
    style=style,
    theme=theme,
    task=task,
    language=language,
    schema_name=schema_name,
    need_search=bool(fused_search),
    fused_search=fused_search,
    forced_file_type=forced_file_type,
)

if ai_raw:
    result_core = normalize_ai_result(
        ai_raw,
        fallback={
            "file_type": forced_file_type or fallback_route["file_type"],
            "intent": fallback_route["intent"],
            "schema_name": schema_name,
            "title": build_title(merged_text),
            "style": style or "默认",
            "content_text": merged_text,
        },
    )
    if forced_file_type:
        result_core["file_type"] = forced_file_type
else:
    result_core = build_fallback_result(
        cleaned=merged_text,
        style=style,
        schema_name=schema_name,
        fallback_route=fallback_route,
        forced_type=forced_file_type,
        fused_search=fused_search,
    )

if not result_core.get("content_text") and result_core.get("file_type") in {"pdf", "docx"}:
    result_core["content_text"] = merged_text

elapsed_ms = int((time.time() - start_ts) * 1000)
output = {
    "ok": True,
    "task_id": task_id,
    "status": "completed",
    "app": APP_NAME,
    "env": APP_ENV,
    "language": language,
    "style": style,
    "theme": theme,
    "need_search": need_search,
    "search": search_data,
    "search_fused": fused_search,
    "result": result_core,
    "filename": sanitize_filename(custom_name or result_core.get("title") or build_title(merged_text)),
    "elapsed_ms": elapsed_ms,
}
update_task(task_id, status="completed", step="done", result_summary={
    "file_type": result_core.get("file_type"),
    "schema_name": result_core.get("schema_name"),
    "elapsed_ms": elapsed_ms,
})
return output
```

# =============================================================================

# PDF

# =============================================================================

def wrap_text_lines(text: str, c, font_name: str, font_size: int, max_width: float):
if not text or not text.strip():
return ["（空内容）"]

```
wrapped = []
raw_lines = text.splitlines()

for raw in raw_lines:
    line = raw.rstrip()
    if not line.strip():
        wrapped.append("")
        continue

    current = ""
    for ch in line:
        test = current + ch
        if c.stringWidth(test, font_name, font_size) <= max_width:
            current = test
        else:
            if current:
                wrapped.append(current)
            current = ch

    if current:
        wrapped.append(current)

return wrapped if wrapped else ["（空内容）"]
```

def create_pdf(text: str, output_path: str, style: str = None, theme: str = None):
pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
palette = get_theme_palette(theme)

```
c = canvas.Canvas(output_path, pagesize=A4)
width, height = A4
title = build_title(text)

c.setFillColorRGB(
    palette["primary"][0] / 255.0,
    palette["primary"][1] / 255.0,
    palette["primary"][2] / 255.0,
)
c.rect(0, height - 90, width, 90, fill=1, stroke=0)

c.setFillColorRGB(1, 1, 1)
c.setFont("STSong-Light", 20)
c.drawCentredString(width / 2, height - 45, title)

c.setFont("STSong-Light", 10)
info = f"风格：{style or '默认'}    主题：{palette['name']}"
c.drawString(40, height - 78, info)

y = height - 120
c.setFillColorRGB(0, 0, 0)
c.setFont("STSong-Light", 12)
max_body_width = width - 80
lines = wrap_text_lines(text, c, "STSong-Light", 12, max_body_width)

for line in lines:
    if y < 50:
        c.showPage()
        y = height - 50
        c.setFont("STSong-Light", 12)
    c.drawString(40, y, line)
    y -= 24

c.save()
```

# =============================================================================

# DOCX

# =============================================================================

def add_doc_divider(doc):
p = doc.add_paragraph()
run = p.add_run("—" * 42)
run.font.size = DocxPt(10)
p.alignment = WD_ALIGN_PARAGRAPH.CENTER

def create_docx(text: str, output_path: str, style: str = None, theme: str = None):
palette = get_theme_palette(theme)
doc = Document()
title = build_title(text)

```
h = doc.add_heading("", level=1)
h.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = h.add_run(title)
run.font.size = DocxPt(20)
run.font.bold = True

meta = doc.add_paragraph()
meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
meta_run = meta.add_run(f"风格：{style or '默认'}    主题：{palette['name']}")
meta_run.font.size = DocxPt(10)

add_doc_divider(doc)
doc.add_paragraph("")

paragraphs = text.splitlines() if text else ["（空内容）"]
for line in paragraphs:
    p = doc.add_paragraph(line if line.strip() else "")
    p.paragraph_format.space_after = DocxPt(8)
    p.paragraph_format.line_spacing = 1.5
    for r in p.runs:
        r.font.size = DocxPt(12)

doc.save(output_path)
```

# =============================================================================

# XLSX

# =============================================================================

def parse_table_text(text: str):
lines = [line.strip() for line in text.splitlines() if line.strip()]
if not lines:
return [["（空内容）"]]

```
rows = []
for line in lines:
    if "\t" in line:
        rows.append([cell.strip() for cell in line.split("\t")])
    elif "," in line:
        rows.append([cell.strip() for cell in line.split(",")])
    elif "，" in line:
        rows.append([cell.strip() for cell in line.split("，")])
    else:
        rows.append([line])
return rows
```

def autofit_worksheet(ws):
for col in ws.columns:
max_length = 0
col_letter = col[0].column_letter
for cell in col:
try:
cell_len = len(str(cell.value)) if cell.value is not None else 0
if cell_len > max_length:
max_length = cell_len
except Exception:
pass
ws.column_dimensions[col_letter].width = min(max_length + 4, 30)

def style_excel_sheet(ws, style: str = None, theme: str = None):
palette = get_theme_palette(theme)
thin = Side(border_style="thin", color="CCCCCC")

```
for row in ws.iter_rows():
    for cell in row:
        cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)
        cell.alignment = Alignment(vertical="center", horizontal="center", wrap_text=True)

if ws.max_row >= 1:
    for cell in ws[1]:
        cell.font = Font(bold=True, color="FFFFFF" if theme == "深色商务" else "000000")
        cell.alignment = Alignment(vertical="center", horizontal="center")
        cell.fill = PatternFill(fill_type="solid", fgColor=palette["excel_header"])
```

def create_xlsx(text: str, output_path: str, style: str = None, theme: str = None):
wb = Workbook()
ws = wb.active
ws.title = "Sheet1"

```
rows = parse_table_text(text)
for row_idx, row in enumerate(rows, start=1):
    for col_idx, value in enumerate(row, start=1):
        ws.cell(row=row_idx, column=col_idx, value=value)

style_excel_sheet(ws, style, theme)
autofit_worksheet(ws)
wb.save(output_path)
```

def create_xlsx_structured(output_path: str, headers, rows, style: str = None, theme: str = None):
wb = Workbook()
ws = wb.active
ws.title = "Sheet1"

```
if not headers:
    headers = ["内容"]
if not rows:
    rows = [["（空内容）"]]

for col_idx, header in enumerate(headers, start=1):
    ws.cell(row=1, column=col_idx, value=header)

for row_idx, row in enumerate(rows, start=2):
    for col_idx, value in enumerate(row, start=1):
        ws.cell(row=row_idx, column=col_idx, value=value)

style_excel_sheet(ws, style, theme)
autofit_worksheet(ws)
wb.save(output_path)
```

# =============================================================================

# PPTX

# =============================================================================

def create_slide_header(slide, title_text, palette):
title_box = slide.shapes.add_textbox(Pt(32), Pt(24), Pt(640), Pt(40))
tf_title = title_box.text_frame
tf_title.clear()
p_title = tf_title.paragraphs[0]
p_title.text = title_text
p_title.font.size = Pt(24)
p_title.font.bold = True
p_title.font.name = "Microsoft YaHei"
p_title.font.color.rgb = palette["primary"]

```
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(32), Pt(70), Pt(640), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = palette["secondary"]
line.line.fill.background()
```

def create_pptx_structured(output_path: str, title: str, slides_data, style: str = None, theme: str = None):
prs = Presentation()
palette = get_theme_palette(theme)

```
slide = prs.slides.add_slide(prs.slide_layouts[6])

block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(48), Pt(110), Pt(620), Pt(120))
block.fill.solid()
block.fill.fore_color.rgb = palette["primary"]
block.line.fill.background()

title_box = slide.shapes.add_textbox(Pt(70), Pt(135), Pt(580), Pt(50))
tf = title_box.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = title
p.font.size = Pt(28)
p.font.bold = True
p.font.name = "Microsoft YaHei"
p.font.color.rgb = RGBColor(255, 255, 255)

subtitle = choose_subtitle_by_style(style)
sub_box = slide.shapes.add_textbox(Pt(70), Pt(185), Pt(580), Pt(30))
tf2 = sub_box.text_frame
tf2.clear()
p2 = tf2.paragraphs[0]
p2.text = f"{subtitle} | 主题：{palette['name']}"
p2.font.size = Pt(14)
p2.font.name = "Microsoft YaHei"
p2.font.color.rgb = RGBColor(255, 255, 255)

if not slides_data:
    slides_data = [{"title": "内容", "bullets": ["（空内容）"]}]

for item in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_header(slide, item.get("title", "内容"), palette)

    content_box = slide.shapes.add_textbox(Pt(52), Pt(110), Pt(610), Pt(330))
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.clear()

    bullets = item.get("bullets", []) or ["（空内容）"]
    bullets = bullets[:4]

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20 if style == "汇报" else 18)
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = palette["text"]
        p.space_after = Pt(10)

    footer = slide.shapes.add_textbox(Pt(52), Pt(470), Pt(300), Pt(20))
    tf_footer = footer.text_frame
    p_footer = tf_footer.paragraphs[0]
    p_footer.text = f"{subtitle}"
    p_footer.font.size = Pt(9)
    p_footer.font.name = "Microsoft YaHei"
    p_footer.font.color.rgb = palette["subtext"]

prs.save(output_path)
```

def create_pptx(text: str, output_path: str, style: str = None, theme: str = None):
title = build_title(text)
lines = [line.strip() for line in text.splitlines() if line.strip()]
bullets = lines[1:] if len(lines) > 1 else ["（空内容）"]
slides_data = [{"title": title, "bullets": bullets[:4]}]
create_pptx_structured(output_path, title, slides_data, style, theme)

# =============================================================================

# File rendering bridge

# =============================================================================

def render_result_to_file(result_payload: Dict[str, Any], custom_name: Optional[str] = None) -> Tuple[str, str]:
result = result_payload["result"]
file_type = result["file_type"]
style = result.get("style")
theme = result_payload.get("theme")
ai_title = result.get("title")
content_text = result.get("content_text", "")
table_headers = result.get("table_headers", [])
table_rows = result.get("table_rows", [])
slides = result.get("slides", [])

```
file_path = build_output_filename(content_text or ai_title or "文件", file_type, custom_name, ai_title)

if file_type == "pdf":
    create_pdf(content_text, file_path, style, theme)
elif file_type == "docx":
    create_docx(content_text, file_path, style, theme)
elif file_type == "xlsx":
    if table_headers or table_rows:
        create_xlsx_structured(file_path, table_headers, table_rows, style, theme)
    else:
        create_xlsx(content_text, file_path, style, theme)
elif file_type == "pptx":
    if slides:
        create_pptx_structured(file_path, ai_title or build_title(content_text), slides, style, theme)
    else:
        create_pptx(content_text, file_path, style, theme)
else:
    raise ValueError(f"unsupported file type: {file_type}")

return file_path, file_type
```

# =============================================================================

# Routes

# =============================================================================

@app.route("/")
def home():
return jsonify({
"ok": True,
"app": APP_NAME,
"env": APP_ENV,
"time": now_iso(),
"routes": [
"GET /health",
"GET /health/detail",
"GET /schemas",
"POST /search",
"POST /generate",
"POST /generate/file",
"GET /tasks/<task_id>",
"POST /webhook",
"GET /set_webhook?token=...",
],
})

@app.route("/health")
def health():
return {"status": "ok"}

@app.route("/health/detail")
def health_detail():
return jsonify({
"status": "ok",
"app": APP_NAME,
"env": APP_ENV,
"bot_token_configured": bool(BOT_TOKEN),
"ai_api_key_configured": bool(AI_API_KEY),
"ai_api_url_configured": bool(AI_API_URL),
"render_external_url_configured": bool(RENDER_EXTERNAL_URL),
"search_provider": SEARCH_PROVIDER,
"search_api_key_configured": bool(SEARCH_API_KEY or SEARCH_PROVIDER != "tavily"),
"task_count": len(TASKS),
"time": now_iso(),
})

@app.route("/schemas")
def schemas():
return jsonify({"ok": True, "schemas": SCHEMAS})

@app.route("/tasks/<task_id>")
def get_task(task_id: str):
task = TASKS.get(task_id)
if not task:
return make_error("task not found", 404)
return jsonify({"ok": True, "task": task})

@app.route("/search", methods=["POST"])
def search_only():
payload = request.get_json(silent=True) or {}
query = clean_text(payload.get("query", ""))
if not query:
return make_error("query is required")

```
try:
    provider = get_search_provider()
    language = payload.get("language") or detect_language(query)
    data = provider.search(
        query=query,
        language=language,
        max_results=int(payload.get("max_results", DEFAULT_MAX_RESULTS)),
        topic=payload.get("topic", DEFAULT_TOPIC),
        search_depth=payload.get("search_depth", DEFAULT_SEARCH_DEPTH),
    )
    fused = fuse_search_results(data, int(payload.get("max_results", DEFAULT_MAX_RESULTS)))
    return jsonify({"ok": True, "search": data, "fused": fused})
except Exception as exc:
    logger.exception("search_only failed")
    return make_error(str(exc), 500)
```

@app.route("/generate", methods=["POST"])
def generate_json_route():
payload = request.get_json(silent=True) or {}
if not payload.get("task") and not payload.get("details") and not payload.get("content"):
return make_error("task or details/content is required")

```
task_id = str(uuid.uuid4())
try:
    result = create_generation_result(payload, task_id)
    return jsonify(result)
except Exception as exc:
    logger.exception("generate_json_route failed")
    update_task(task_id, status="failed", step="error", error=str(exc))
    return make_error(str(exc), 500, extra={"task_id": task_id})
```

@app.route("/generate/file", methods=["POST"])
def generate_file_route():
payload = request.get_json(silent=True) or {}
if not payload.get("task") and not payload.get("details") and not payload.get("content"):
return make_error("task or details/content is required")

```
task_id = str(uuid.uuid4())
file_path = None
try:
    result = create_generation_result(payload, task_id)
    file_path, file_type = render_result_to_file(result, payload.get("filename"))
    update_task(task_id, step="file_rendered", output_file=os.path.basename(file_path))

    mime_map = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    return send_file(
        file_path,
        mimetype=mime_map[file_type],
        as_attachment=True,
        download_name=os.path.basename(file_path),
    )
except Exception as exc:
    logger.exception("generate_file_route failed")
    update_task(task_id, status="failed", step="error", error=str(exc))
    return make_error(str(exc), 500, extra={"task_id": task_id})
finally:
    try:
        if file_path and os.path.exists(file_path):
            os.remove(file_path)
    except Exception:
        logger.exception("failed to cleanup temp file")
```

@app.route("/webhook", methods=["POST"])
def webhook():
secret = request.headers.get("X-Telegram-Bot-Api-Secret-Token")
if secret != WEBHOOK_SECRET:
return jsonify({"ok": False, "error": "unauthorized"}), 403

```
update = request.get_json(silent=True) or {}
message = update.get("message", {})
chat = message.get("chat", {})
chat_id = chat.get("id")
text = message.get("text", "")
reply_to_message = message.get("reply_to_message", {})

logger.info("incoming telegram message chat_id=%s text=%s", chat_id, (text or "")[:200])

if not chat_id:
    return jsonify({"ok": True})

if text == "/start":
    send_message(
        chat_id,
        "你好，直接发送内容即可。我会自动识别并生成 PDF / Word / Excel / PPT 文件。支持文件名、风格、主题，也支持回复上一条消息后补一句发我pdf。",
    )
    return jsonify({"ok": True})

if text.startswith("/") and text != "/start":
    send_message(chat_id, "请直接发送正文内容，不需要输入命令。")
    return jsonify({"ok": True})

merged_text = extract_effective_request(text, reply_to_message)
custom_name, text_without_filename = extract_custom_filename(merged_text)
style, text_without_style = extract_style(text_without_filename)
theme, text_without_theme = extract_theme(text_without_style)
cleaned = strip_control_phrases(text_without_theme)

if not cleaned:
    send_message(chat_id, "请发送你要生成文件的文字内容。")
    return jsonify({"ok": True})

if len(cleaned) > MAX_INPUT_CHARS:
    send_message(chat_id, "内容过长，请拆分后再发送。")
    return jsonify({"ok": True})

lowered = cleaned.lower()
forced_type = None
if "pdf" in lowered:
    forced_type = "pdf"
elif "docx" in lowered or "word" in lowered:
    forced_type = "docx"
elif "xlsx" in lowered or "excel" in lowered:
    forced_type = "xlsx"
elif "pptx" in lowered or re.search(r"\bppt\b", lowered):
    forced_type = "pptx"

task_id = str(uuid.uuid4())
file_path = None

try:
    send_message(chat_id, "已收到，正在为你生成文件...")
    result = create_generation_result(
        {
            "task": build_title(cleaned),
            "content": cleaned,
            "style": style or "默认",
            "theme": theme or "商务蓝",
            "filename": custom_name,
            "format": forced_type,
        },
        task_id,
    )

    file_path, file_type = render_result_to_file(result, custom_name)

    caption_map = {
        "pdf": "你的 PDF 已生成",
        "docx": "你的 Word 文档已生成",
        "xlsx": "你的 Excel 文件已生成",
        "pptx": "你的 PPT 文件已生成",
    }
    send_document(chat_id, file_path, caption_map.get(file_type, "文件已生成"))
    return jsonify({"ok": True})
except Exception as e:
    logger.exception("telegram generate failed")
    try:
        send_message(chat_id, f"生成失败：{str(e)}")
    except Exception:
        logger.exception("failed to send telegram error message")
    return jsonify({"ok": False, "error": str(e)}), 500
finally:
    try:
        if file_path and os.path.exists(file_path):
            os.remove(file_path)
    except Exception:
        logger.exception("failed to cleanup telegram temp file")
```

@app.route("/set_webhook")
def set_webhook():
token = request.args.get("token", "")
if not ADMIN_TOKEN or token != ADMIN_TOKEN:
return {"ok": False, "error": "unauthorized"}, 403

```
url = RENDER_EXTERNAL_URL.rstrip("/") + "/webhook"
if not BOT_TOKEN or not url:
    return {
        "ok": False,
        "error": "missing BOT_TOKEN or RENDER_EXTERNAL_URL",
    }, 400

resp = requests.post(
    f"{TELEGRAM_API}/setWebhook",
    json={
        "url": url,
        "secret_token": WEBHOOK_SECRET,
    },
    timeout=30,
)
return resp.json(), resp.status_code
```

if **name** == "**main**":
print(
f"{APP_NAME} started. host={HOST} port={PORT} env={APP_ENV} "
f"bot={'yes' if BOT_TOKEN else 'no'} ai={'yes' if AI_API_KEY else 'no'} search={SEARCH_PROVIDER}"
)
app.run(host=HOST, port=PORT, debug=DEBUG)
